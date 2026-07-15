"""
Client Reports Automation - Microservice v2
============================================

Flask microservice that generates .pptx reports from templates by:
1. Replacing {{placeholder}} text tags with actual data
2. Replacing images on specified slides with chart images from URLs
   (preserving aspect ratio, centering within the target box)

Endpoints:
- GET  /          : health check
- POST /generate  : generate a report from JSON body

Deployed on Render.com free tier.
"""
from flask import Flask, request, send_file, jsonify
from pptx import Presentation
import requests
import io
import os
import traceback

app = Flask(__name__)

TEMPLATES_DIR = os.path.dirname(os.path.abspath(__file__))
REQUEST_TIMEOUT = 30


# ==========================================
# TEXT REPLACEMENT
# ==========================================

def replace_placeholders_in_paragraph(paragraph, placeholders):
    """Handle placeholders that span multiple runs by consolidating them."""
    full_text = ''.join(run.text for run in paragraph.runs)
    if '{{' not in full_text:
        return False

    new_text = full_text
    for key, value in placeholders.items():
        marker = '{{' + key + '}}'
        if marker in new_text:
            replacement = str(value) if value is not None else ''
            new_text = new_text.replace(marker, replacement)

    if new_text == full_text:
        return False

    if paragraph.runs:
        paragraph.runs[0].text = new_text
        for run in paragraph.runs[1:]:
            run.text = ''

    return True


def replace_in_shape(shape, placeholders):
    """Recursively process a shape for text replacement."""
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            replace_placeholders_in_paragraph(para, placeholders)

    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    replace_placeholders_in_paragraph(para, placeholders)

    try:
        if shape.shape_type == 6:  # GROUP
            for sub_shape in shape.shapes:
                replace_in_shape(sub_shape, placeholders)
    except Exception:
        pass


# ==========================================
# IMAGE REPLACEMENT (with aspect ratio preservation)
# ==========================================

def get_image_dimensions(image_bytes):
    """Get (width, height) in pixels from image bytes."""
    from PIL import Image
    try:
        img = Image.open(io.BytesIO(image_bytes))
        return img.size
    except Exception:
        return None


def replace_image_on_slide(slide, image_bytes, image_index=0):
    """
    Replace the Nth largest picture on a slide.
    Preserves the new image's aspect ratio and centers it within the target box.
    """
    pictures = [s for s in slide.shapes if s.shape_type == 13]
    if not pictures:
        return False

    pictures.sort(key=lambda p: (p.width or 0) * (p.height or 0), reverse=True)

    if image_index >= len(pictures):
        return False

    target = pictures[image_index]
    box_left, box_top = target.left, target.top
    box_width, box_height = target.width, target.height

    sp = target._element
    sp.getparent().remove(sp)

    dims = get_image_dimensions(image_bytes)
    if dims is None:
        slide.shapes.add_picture(io.BytesIO(image_bytes), box_left, box_top, box_width, box_height)
        return True

    img_w_px, img_h_px = dims
    img_ratio = img_w_px / img_h_px if img_h_px > 0 else 1
    box_ratio = box_width / box_height if box_height > 0 else 1

    if img_ratio > box_ratio:
        new_width = box_width
        new_height = int(box_width / img_ratio)
    else:
        new_height = box_height
        new_width = int(box_height * img_ratio)

    new_left = box_left + (box_width - new_width) // 2
    new_top = box_top + (box_height - new_height) // 2

    slide.shapes.add_picture(io.BytesIO(image_bytes), new_left, new_top, new_width, new_height)
    return True


# ==========================================
# ENDPOINTS
# ==========================================

@app.route('/', methods=['GET'])
def health():
    templates = [f for f in os.listdir(TEMPLATES_DIR) if f.endswith('.pptx')]
    return jsonify({
        'status': 'ok',
        'service': 'client-reports-automation',
        'templates_dir': TEMPLATES_DIR,
        'available_templates': templates
    })


@app.route('/generate', methods=['POST'])
def generate_report():
    try:
        data = request.get_json(force=True, silent=False)
        if not data:
            return jsonify({'error': 'Empty or invalid JSON body'}), 400

        template_name = data.get('template', 'Malne-Template.pptx')
        placeholders = data.get('placeholders', {})
        images = data.get('images', {})

        template_path = os.path.join(TEMPLATES_DIR, template_name)
        if not os.path.exists(template_path):
            available = [f for f in os.listdir(TEMPLATES_DIR) if f.endswith('.pptx')]
            return jsonify({
                'error': f'Template not found: {template_name}',
                'available': available
            }), 404

        prs = Presentation(template_path)

        for slide in prs.slides:
            for shape in slide.shapes:
                replace_in_shape(shape, placeholders)

        images_replaced = 0
        for slide_num_str, image_url in images.items():
            try:
                slide_num = int(slide_num_str)
                slide_idx = slide_num - 1
                if slide_idx < 0 or slide_idx >= len(prs.slides):
                    print(f'[WARN] Slide {slide_num} out of range')
                    continue

                print(f'[INFO] Downloading image for slide {slide_num}...')
                resp = requests.get(image_url, timeout=REQUEST_TIMEOUT)
                if resp.status_code != 200:
                    print(f'[WARN] Image download failed slide {slide_num}: HTTP {resp.status_code}')
                    continue

                slide = prs.slides[slide_idx]
                if replace_image_on_slide(slide, resp.content):
                    images_replaced += 1
                    print(f'[OK] Image replaced on slide {slide_num}')
                else:
                    print(f'[WARN] No picture found on slide {slide_num} to replace')
            except Exception as e:
                print(f'[ERROR] Image replacement slide {slide_num_str}: {e}')
                continue

        output = io.BytesIO()
        prs.save(output)
        output.seek(0)

        template_base = os.path.splitext(template_name)[0]
        filename = f'{template_base}-generated.pptx'

        print(f'[SUCCESS] Report generated. Images replaced: {images_replaced}')

        return send_file(
            output,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name=filename
        )

    except Exception as e:
        traceback.print_exc()
        return jsonify({
            'error': str(e),
            'traceback': traceback.format_exc().split('\n')[-10:]
        }), 500


if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5000))
    app.run(host='0.0.0.0', port=port, debug=False)
